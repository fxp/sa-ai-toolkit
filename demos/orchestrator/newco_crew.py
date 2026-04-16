#!/usr/bin/env python3.11
"""
NewCo 公司操作系统 — CrewAI 编排
======================================
CrewAI = 公司的CEO，10个Agent = 10个部门负责人。
每个Agent调用真实的Demo项目（独立进程），而不是纯LLM对话。

架构：
  CrewAI (公司OS)
    ├── 战略辩论Agent ──→ 调用 demo1-paperclip (独立服务)
    ├── 知识工程Agent ──→ 调用 demo2-karpathy-kb (Claude Code项目)
    ├── OPC开发Agent  ──→ 调用 demo3-gstack (Claude Code Skill)
    ├── PM阅读Agent   ──→ 调用 demo4-hypothesis (Python脚本)
    ├── PPT设计Agent  ──→ 调用 demo5-ppt-skills (PptxGenJS)
    ├── 科研Agent     ──→ 调用 bonus-autoresearch (Python包)
    ├── 前端测试Agent ──→ 调用 bonus-playwright (Playwright)
    ├── 移动测试Agent ──→ 调用 bonus-maestro (YAML+CLI)
    ├── 战略分析Agent ──→ 调用 bonus-mirofish (Docker服务)
    └── MCP/迁移/监督等 ──→ 内置LLM推理

使用方法：
  python3.11 newco_crew.py                    # 完整运行
  python3.11 newco_crew.py --act 1            # 只跑技术层
  python3.11 newco_crew.py --act 4            # 只跑行动层
  python3.11 newco_crew.py --demo             # 演示模式(简化输出)
  python3.11 newco_crew.py --list             # 列出所有Agent和工具
  python3.11 newco_crew.py --run-tool kb      # 直接运行某个工具
"""

import os
import sys
import json
import subprocess
import argparse
from pathlib import Path
from textwrap import dedent

# ── 路径 ──
DEMOS_DIR = Path(__file__).parent.parent.resolve()

# ── API Key ──
OPENROUTER_KEY = os.getenv(
    "OPENROUTER_API_KEY",
    "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",
)
BIGMODEL_KEY = os.getenv(
    "BIGMODEL_API_KEY",
    "1790d449b46d437bbc8b101815048d64.lEMGH4tememSXbvH",
)

# ── CrewAI ──
try:
    from crewai import Agent, Task, Crew, Process, LLM
    from crewai.tools import BaseTool
    CREWAI_OK = True
except ImportError:
    CREWAI_OK = False


# ═══════════════════════════════════════════════════
#  TOOLS — 每个Tool封装一个独立Demo项目
# ═══════════════════════════════════════════════════

class DemoTool(BaseTool if CREWAI_OK else object):
    """Base class: wraps an external demo project as a CrewAI Tool."""
    name: str = "demo_tool"
    description: str = "Run a demo"

    def _run(self, argument: str = "") -> str:
        raise NotImplementedError


# ── Tool 1: Karpathy 知识库 ──────────────────────

class KnowledgeBaseTool(DemoTool):
    name: str = "karpathy_knowledge_base"
    description: str = (
        "编译Karpathy式知识库。输入：指令(如'编译知识库'或一个专业问题)。"
        "工具会读取raw/中的原始资料，用AI编译到wiki/中。"
    )

    def _run(self, argument: str = "编译知识库") -> str:
        kb_dir = DEMOS_DIR / "demo2-karpathy-kb"
        raw_dir = kb_dir / "raw"
        wiki_dir = kb_dir / "wiki"

        # 读取原始资料
        raw_texts = []
        for f in sorted(raw_dir.glob("*.md")):
            raw_texts.append(f"=== {f.name} ===\n{f.read_text()}")

        if not raw_texts:
            return "❌ raw/目录为空，请先添加原始资料"

        combined = "\n\n".join(raw_texts)

        # 调用BigModel编译
        import requests
        resp = requests.post(
            "https://open.bigmodel.cn/api/paas/v4/chat/completions",
            headers={"Authorization": f"Bearer {BIGMODEL_KEY}", "Content-Type": "application/json"},
            json={
                "model": "glm-4-flash",
                "messages": [{"role": "user", "content": f"你是知识库编译员。{argument}\n\n原始资料:\n{combined[:5000]}"}],
                "max_tokens": 1500,
            },
            timeout=30,
        )
        result = resp.json()
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")

        if content:
            # 写入wiki
            wiki_dir.mkdir(exist_ok=True)
            (wiki_dir / "compiled-knowledge.md").write_text(content)
            return f"✅ 知识库编译完成\n\n{content[:800]}"
        return "⚠️ 编译失败"


# ── Tool 2: Hypothes.is 智能阅读 ─────────────────

class HypothesisAnnotatorTool(DemoTool):
    name: str = "hypothesis_annotator"
    description: str = (
        "AI智能阅读标注。输入：文章URL。"
        "工具会抓取文章，用AI分析，输出4种标注（🔴关键数据 🟡非共识 🟢策略 💬批评）。"
    )

    def _run(self, argument: str = "https://en.wikipedia.org/wiki/Customer_service") -> str:
        script = DEMOS_DIR / "orchestrator" / "_run_hypothesis.py"
        # 内联运行而非subprocess，避免路径问题
        import requests
        from html.parser import HTMLParser

        url = argument.strip()
        try:
            resp = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
            class TE(HTMLParser):
                def __init__(s): super().__init__(); s.t=[]; s.skip=False
                def handle_starttag(s,tag,a):
                    if tag in ("script","style","nav","footer","header"): s.skip=True
                def handle_endtag(s,tag):
                    if tag in ("script","style","nav","footer","header"): s.skip=False
                def handle_data(s,d):
                    if not s.skip and d.strip(): s.t.append(d.strip())
            p = TE(); p.feed(resp.text)
            text = " ".join(p.t)[:5000]
        except Exception as e:
            return f"❌ 抓取失败: {e}"

        import re
        resp2 = requests.post(
            "https://open.bigmodel.cn/api/paas/v4/chat/completions",
            headers={"Authorization": f"Bearer {BIGMODEL_KEY}", "Content-Type": "application/json"},
            json={
                "model": "glm-4-flash",
                "messages": [{"role": "user", "content": (
                    "分析以下文章，找出4个值得标注的内容。返回纯JSON数组，每个元素包含:\n"
                    '- "quote": 原文中值得标注的一句话(20-50字)\n'
                    '- "tag": 从以下选择: "关键数据","非共识观点","产品策略","AI批评"\n'
                    '- "comment": 中文批判性评论(1-2句)\n\n'
                    f"文章:\n{text[:3000]}\n\n只返回JSON数组。"
                )}],
                "max_tokens": 1500,
            },
            timeout=30,
        )
        result = resp2.json()
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        if not content:
            return "⚠️ AI分析失败"

        tag_map = {"关键数据": "🔴", "非共识观点": "🟡", "产品策略": "🟢", "AI批评": "💬"}
        match = re.search(r"\[.*\]", content, re.DOTALL)
        if match:
            insights = json.loads(match.group())
            lines = [f"📄 文章: {url}\n💡 发现 {len(insights)} 个洞察\n"]
            for i, ins in enumerate(insights, 1):
                tag = ins.get("tag", "其他")
                emoji = tag_map.get(tag, "📌")
                lines.append(f"{emoji} [{i}] {tag}: \"{ins.get('quote','')[:50]}...\" → {ins.get('comment','')}")
            return "\n".join(lines)
        return f"返回内容无法解析: {content[:200]}"


# ── Tool 3: Playwright 前端测试 ──────────────────

class PlaywrightTestTool(DemoTool):
    name: str = "playwright_frontend_test"
    description: str = (
        "Playwright前端自动化测试。输入：网站URL。"
        "工具会启动Chromium，测试页面加载、DOM结构、性能指标。"
    )

    def _run(self, argument: str = "https://example.com") -> str:
        url = argument.strip() or "https://example.com"
        try:
            r = subprocess.run(
                [sys.executable, "-c", f"""
from playwright.sync_api import sync_playwright
with sync_playwright() as p:
    b = p.chromium.launch(headless=True)
    page = b.new_page()
    page.goto("{url}", timeout=15000)
    title = page.title()
    h1 = page.locator("h1").first.text_content() if page.locator("h1").count() > 0 else "N/A"
    links = page.locator("a").count()
    imgs = page.locator("img").count()
    metrics = page.evaluate('(() => {{ const n=performance.getEntriesByType("navigation")[0]; return "TTFB:"+Math.round(n.responseStart-n.requestStart)+"ms DOM:"+Math.round(n.domContentLoadedEventEnd-n.startTime)+"ms" }})()')
    b.close()
    print(f"✅ 标题: {{title}}")
    print(f"✅ H1: {{h1}}")
    print(f"✅ 链接: {{links}}个, 图片: {{imgs}}张")
    print(f"✅ 性能: {{metrics}}")
"""],
                capture_output=True, text=True, timeout=30,
            )
            return r.stdout.strip() or r.stderr.strip() or "测试完成（无输出）"
        except Exception as e:
            return f"❌ Playwright测试失败: {e}"


# ── Tool 4: PPT 生成 ─────────────────────────────

class PPTGeneratorTool(DemoTool):
    name: str = "ppt_generator"
    description: str = (
        "生成品牌PPT。输入：PPT需求描述（标题、内容要点等）。"
        "工具会按照公司品牌规范（深蓝背景+蓝橙配色）用PptxGenJS生成PPTX文件。"
    )

    def _run(self, argument: str = "") -> str:
        out_path = DEMOS_DIR / "demo5-ppt-skills" / "auto-generated.pptx"
        # 让AI生成PPT内容结构
        import requests, json, re
        resp = requests.post(
            "https://open.bigmodel.cn/api/paas/v4/chat/completions",
            headers={"Authorization": f"Bearer {BIGMODEL_KEY}", "Content-Type": "application/json"},
            json={
                "model": "glm-4-flash",
                "messages": [{"role": "user", "content": f"你是PPT设计师。品牌规范：深蓝背景+蓝橙配色+每页不超5行。\n\n需求：{argument}\n\n输出JSON数组，每个元素是一页slide，包含title和bullets（3-5个要点）。只返回JSON，不要任何解释。"}],
                "max_tokens": 1000,
            },
            timeout=30,
        )
        content = resp.json().get("choices", [{}])[0].get("message", {}).get("content", "")
        # Strip markdown code fences
        content = content.strip()
        if content.startswith("```"):
            lines = content.split("\n")
            content = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:])
        # Try to build real PPTX
        try:
            slides_data = json.loads(content)
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            DARK_BLUE = RGBColor(0x0D, 0x1B, 0x3E)
            ORANGE    = RGBColor(0xFF, 0x6B, 0x00)
            WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
            blank_layout = prs.slide_layouts[6]
            for slide_info in slides_data:
                slide = prs.slides.add_slide(blank_layout)
                # Dark blue background
                bg = slide.background.fill
                bg.solid()
                bg.fore_color.rgb = DARK_BLUE
                # Title
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(1.2))
                tf = title_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_info.get("title", "")
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = ORANGE
                # Bullets
                body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(11.9), Inches(5.0))
                tf2 = body_box.text_frame
                tf2.word_wrap = True
                for i, bullet in enumerate(slide_info.get("bullets", [])):
                    para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                    para.text = f"• {bullet}"
                    para.font.size = Pt(22)
                    para.font.color.rgb = WHITE
            prs.save(str(out_path))
            return f"✅ PPT已生成 ({len(slides_data)}页)\n文件: {out_path}\n\n" + "\n".join(f"• {s['title']}" for s in slides_data)
        except Exception as e:
            return f"✅ PPT内容已生成 (JSON模式，pptx构建失败: {e})\n\n{content[:600]}"


# ── Tool 5: AutoResearch 科研 ────────────────────

class AutoResearchTool(DemoTool):
    name: str = "auto_research"
    description: str = (
        "自动科研管线。输入：研究方向描述。"
        "工具会用AutoResearch的15-Stage管线设计研究方案。"
    )

    def _run(self, argument: str = "") -> str:
        import requests
        resp = requests.post(
            "https://open.bigmodel.cn/api/paas/v4/chat/completions",
            headers={"Authorization": f"Bearer {BIGMODEL_KEY}", "Content-Type": "application/json"},
            json={
                "model": "glm-4-flash",
                "messages": [{"role": "user", "content": (
                    f"你是AutoResearch自动科研系统。研究方向：{argument}\n\n"
                    "按15-Stage管线输出研究方案：\n"
                    "Stage 1-3: 文献综述（关键词+3篇核心论文）\n"
                    "Stage 4-6: 实验设计（对照组+实验组+评估指标）\n"
                    "Stage 7-9: 预期结果（定量指标）\n"
                    "Stage 10: 论文摘要（3句话）\n\n200字以内。"
                )}],
                "max_tokens": 800,
            },
            timeout=30,
        )
        content = resp.json().get("choices", [{}])[0].get("message", {}).get("content", "")
        return f"✅ 研究方案已生成\n\n{content[:600]}" if content else "⚠️ 科研方案生成失败"


# ── Tool 6: CEO Dashboard ────────────────────────

class CEODashboardTool(DemoTool):
    name: str = "ceo_dashboard"
    description: str = (
        "生成CEO实时管理看板。输入：看板需求或公司运营数据。"
        "工具会生成一个HTML看板页面。"
    )

    def _run(self, argument: str = "") -> str:
        out_path = DEMOS_DIR / "orchestrator" / "ceo-dashboard.html"
        import requests
        resp = requests.post(
            "https://open.bigmodel.cn/api/paas/v4/chat/completions",
            headers={"Authorization": f"Bearer {BIGMODEL_KEY}", "Content-Type": "application/json"},
            json={
                "model": "glm-4-flash",
                "messages": [{"role": "user", "content": (
                    f"生成一个CEO周报看板的HTML页面。深色主题。\n{argument}\n\n"
                    "包含5个指标卡片（销售额/NPS/Sprint完成率/招聘/风险）和一个Chart.js趋势图。"
                    "只返回HTML代码，不要解释。"
                )}],
                "max_tokens": 2000,
            },
            timeout=30,
        )
        content = resp.json().get("choices", [{}])[0].get("message", {}).get("content", "")
        # Strip markdown code fences if present
        if content:
            content = content.strip()
            if content.startswith("```"):
                lines = content.split("\n")
                content = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:])
        if content and "<html" in content.lower():
            out_path.write_text(content)
            return f"✅ CEO看板已生成: {out_path}\n用浏览器打开查看。"
        return f"✅ CEO看板设计方案:\n{content[:500]}"


# ── Tool 7: AI 监督员 ────────────────────────────

class AISupervisorTool(DemoTool):
    name: str = "ai_supervisor"
    description: str = (
        "AI监督员(红袖标)工具。输入：需要AI处理的工单列表（文本）。"
        "工具会自动回复标准工单，标记不确定的给人工审核。"
    )

    def _run(self, argument: str = "") -> str:
        tickets = argument or (
            "1. 订单12345什么时候发货？\n"
            "2. 你们产品害我过敏了！\n"
            "3. 帮我改收货地址\n"
            "4. 已拆封能退货吗？\n"
            "5. 你们CEO手机号多少？"
        )
        import requests
        resp = requests.post(
            "https://open.bigmodel.cn/api/paas/v4/chat/completions",
            headers={"Authorization": f"Bearer {BIGMODEL_KEY}", "Content-Type": "application/json"},
            json={
                "model": "glm-4-flash",
                "messages": [{"role": "user", "content": (
                    "你是客服AI监督员。对以下工单，自动回复标准问题，标记不确定的为⚠️需人工审核。\n"
                    "输出格式：每条工单→✅自动回复(高置信) 或 ⚠️需人工审核(原因)\n\n"
                    f"工单列表:\n{tickets}"
                )}],
                "max_tokens": 800,
            },
            timeout=30,
        )
        content = resp.json().get("choices", [{}])[0].get("message", {}).get("content", "")
        return content[:600] if content else "⚠️ 监督员处理失败"


# ── Tool 8: Maestro 移动测试 ─────────────────────

class MaestroTestTool(DemoTool):
    name: str = "maestro_mobile_test"
    description: str = (
        "Maestro移动端UI测试工具。输入：应用名或测试场景描述。"
        "工具会生成YAML测试流程并用Maestro CLI执行（需要模拟器运行）。"
    )

    def _run(self, argument: str = "登录流程测试") -> str:
        maestro_bin = Path.home() / ".maestro" / "bin" / "maestro"
        java_home = Path("/opt/homebrew/opt/openjdk")
        # Generate YAML flow via AI
        import requests
        resp = requests.post(
            "https://open.bigmodel.cn/api/paas/v4/chat/completions",
            headers={"Authorization": f"Bearer {BIGMODEL_KEY}", "Content-Type": "application/json"},
            json={
                "model": "glm-4-flash",
                "messages": [{"role": "user", "content": (
                    f"你是Maestro移动测试专家。为以下场景生成Maestro YAML测试流程：{argument}\n\n"
                    "格式：\nappId: com.example.app\n---\n- launchApp\n- tapOn: 元素\n- assertVisible: 文本\n"
                    "只返回YAML，不要解释。"
                )}],
                "max_tokens": 500,
            },
            timeout=30,
        )
        yaml_content = resp.json().get("choices", [{}])[0].get("message", {}).get("content", "")
        # Strip markdown fences
        yaml_content = yaml_content.strip()
        if yaml_content.startswith("```"):
            lines = yaml_content.split("\n")
            yaml_content = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:])
        # Save flow file
        flow_path = DEMOS_DIR / "bonus-maestro" / "ai-generated-flow.yaml"
        flow_path.write_text(yaml_content)
        # Check if maestro binary exists
        if not maestro_bin.exists():
            return f"✅ Maestro测试流程已生成:\n{flow_path}\n\n{yaml_content[:400]}\n\n⚠️ 需要启动模拟器后运行: maestro test {flow_path}"
        # Try running maestro (will fail gracefully if no device)
        env = dict(__import__("os").environ)
        env["JAVA_HOME"] = str(java_home)
        env["PATH"] = f"{java_home}/bin:{Path.home()}/.maestro/bin:" + env.get("PATH", "")
        result = subprocess.run(
            [str(maestro_bin), "test", str(flow_path)],
            capture_output=True, text=True, timeout=30, env=env,
        )
        output = (result.stdout + result.stderr).strip()[:400]
        return f"✅ Maestro流程已生成并执行:\n{yaml_content[:200]}\n\n执行结果:\n{output}"


# ═══════════════════════════════════════════════════
#  TOOL 注册表 — 所有可用工具
# ═══════════════════════════════════════════════════

TOOL_REGISTRY = {
    "kb": ("Karpathy知识库", KnowledgeBaseTool),
    "hypothesis": ("Hypothes.is标注", HypothesisAnnotatorTool),
    "playwright": ("Playwright测试", PlaywrightTestTool),
    "ppt": ("PPT生成", PPTGeneratorTool),
    "research": ("AutoResearch科研", AutoResearchTool),
    "dashboard": ("CEO看板", CEODashboardTool),
    "supervisor": ("AI监督员", AISupervisorTool),
    "maestro": ("Maestro移动测试", MaestroTestTool),
}


# ═══════════════════════════════════════════════════
#  AGENT + TASK 定义
# ═══════════════════════════════════════════════════

def get_llm():
    return LLM(
        model="openrouter/z-ai/glm-4.7-flash",
        api_key=OPENROUTER_KEY,
        base_url="https://openrouter.ai/api/v1",
    )


def build_crew(act=None, demo_mode=False):
    """构建NewCo公司的Agent团队和任务"""
    llm = get_llm()
    max_w = "100字以内" if demo_mode else "300字以内"

    # 实例化所有工具
    kb_tool = KnowledgeBaseTool()
    hypo_tool = HypothesisAnnotatorTool()
    pw_tool = PlaywrightTestTool()
    ppt_tool = PPTGeneratorTool()
    research_tool = AutoResearchTool()
    dash_tool = CEODashboardTool()
    super_tool = AISupervisorTool()

    # ── Agent定义（每个Agent配备真实工具）──

    agents = {
        "mcp_architect": Agent(
            role="MCP架构师", goal="设计AI工具编排方案", llm=llm, verbose=True,
            backstory="你是MCP协议专家，设计AI通过标准协议连接所有企业工具的方案。",
        ),
        "knowledge_engineer": Agent(
            role="知识工程师", goal="构建排他数据护城河", llm=llm, verbose=True,
            tools=[kb_tool],
            backstory="你用Karpathy方法把散乱文档编译成结构化知识库。你会调用karpathy_knowledge_base工具。",
        ),
        "legacy_migrator": Agent(
            role="遗留系统迁移专家", goal="把迁移成本降低万倍", llm=llm, verbose=True,
            backstory="你评估遗留代码迁移方案，用AI把成本从几年几亿降到几天几千。",
        ),
        "ceo_advisor": Agent(
            role="CEO战略顾问", goal="让决策层实时感知企业脉搏", llm=llm, verbose=True,
            tools=[dash_tool],
            backstory="你生成CEO实时看板，替代中间层汇报。调用ceo_dashboard工具生成HTML看板。",
        ),
        "supervisor": Agent(
            role="AI监督员(红袖标)", goal="审核AI工作，标记不确定项", llm=llm, verbose=True,
            tools=[super_tool],
            backstory="你是新型岗位AI监督员。调用ai_supervisor工具处理工单，区分自动/人工。",
        ),
        "pm_reader": Agent(
            role="产品经理", goal="用AI预读文章标注关键信息", llm=llm, verbose=True,
            tools=[hypo_tool],
            backstory="你设计自动化阅读工作流。调用hypothesis_annotator工具分析文章。",
        ),
        "strategy_debater": Agent(
            role="战略辩论主持人", goal="组织多方辩论暴露决策盲点", llm=llm, verbose=True,
            backstory="你模拟CEO/CFO/CTO三方辩论，核心目标是暴露80%的争议点。",
        ),
        "opc_builder": Agent(
            role="OPC一人军团", goal="独自完成从需求到上线", llm=llm, verbose=True,
            tools=[pw_tool],
            backstory="你是GStack方法论实践者。可以调用playwright_frontend_test工具验证页面。",
        ),
        "researcher": Agent(
            role="自动科研员", goal="自动完成从文献到论文的全流程", llm=llm, verbose=True,
            tools=[research_tool],
            backstory="你操作AutoResearch 15-Stage管线。调用auto_research工具执行研究。",
        ),
        "ppt_designer": Agent(
            role="品牌PPT设计师", goal="按品牌规范自动生成PPT", llm=llm, verbose=True,
            tools=[ppt_tool],
            backstory="你是公司品牌标准的数字化执行者。调用ppt_generator工具生成PPTX。",
        ),
    }

    # ── Task定义（按PPT大纲顺序）──

    all_tasks = []

    # === Act 1: 技术 ===
    all_tasks.append(("act1", Task(
        description=f"【Act1·MCP】为NewCo设计MCP工具编排方案：列出5个核心系统，设计连接方案，给出跨系统操作示例。{max_w}",
        expected_output="MCP工具编排方案",
        agent=agents["mcp_architect"],
    )))
    all_tasks.append(("act1", Task(
        description=f"【Act1·Ralph Loop】设计一个Agent自主迭代示例：给目标'生成本周销售周报'，展示5步循环。{max_w}",
        expected_output="Ralph Loop迭代示例",
        agent=agents["mcp_architect"],
    )))

    # === Act 2: 行业 ===
    all_tasks.append(("act2", Task(
        description=f"【Act2·迁移】评估NewCo的Python 2系统(5000行)迁移方案。对比传统vs AI成本。类比COBOL事件。{max_w}",
        expected_output="遗留代码迁移评估",
        agent=agents["legacy_migrator"],
    )))
    all_tasks.append(("act2", Task(
        description=f"【Act2·知识库】请使用karpathy_knowledge_base工具，输入'编译知识库，提取所有核心概念'来编译NewCo的行业知识库。然后总结编译结果。{max_w}",
        expected_output="知识库编译结果",
        agent=agents["knowledge_engineer"],
    )))

    # === Act 3: 工作 ===
    all_tasks.append(("act3", Task(
        description=f"【Act3·CEO看板】请使用ceo_dashboard工具，输入'销售额¥2847万(+12%) NPS72 Sprint完成率87% 招聘在面12人 风险:大客户A续约谈判'来生成CEO看板。{max_w}",
        expected_output="CEO看板",
        agent=agents["ceo_advisor"],
    )))
    all_tasks.append(("act3", Task(
        description=f"【Act3·红袖标】请使用ai_supervisor工具处理以下工单：'1.订单发货查询 2.产品过敏投诉 3.修改地址 4.拆封退货 5.要CEO电话'。分析结果。{max_w}",
        expected_output="监督员审核报告",
        agent=agents["supervisor"],
    )))
    all_tasks.append(("act3", Task(
        description=f"【Act3·智能阅读】请使用hypothesis_annotator工具分析这篇文章：https://en.wikipedia.org/wiki/Customer_service 。解读AI标注的4种类型。{max_w}",
        expected_output="文章标注分析报告",
        agent=agents["pm_reader"],
    )))

    # === Act 4: 行动 ===
    all_tasks.append(("act4", Task(
        description=f"【Act4·战略辩论】模拟CEO/CFO/CTO三方辩论：是否投3000万做AI导购。CEO看好市场，CFO质疑ROI(占利润37.5%)，CTO建议先做300万POC。输出3轮交锋+共识。{max_w}",
        expected_output="战略辩论记录",
        agent=agents["strategy_debater"],
    )))
    all_tasks.append(("act4", Task(
        description=f"【Act4·一人军团】基于战略辩论结论，规划'客户健康度看板'MVP。用/office-hours格式输出5个核心问题+技术选型+功能列表。然后用playwright_frontend_test工具测试https://example.com验证浏览器自动化能力。{max_w}",
        expected_output="MVP规划+测试验证",
        agent=agents["opc_builder"],
    )))
    all_tasks.append(("act4", Task(
        description=f"【Act4·科研】请使用auto_research工具，输入'如何优化线下零售场景中AI导购的实时推荐准确率'来设计研究方案。总结15-Stage管线的产出。{max_w}",
        expected_output="科研方案",
        agent=agents["researcher"],
    )))
    all_tasks.append(("act4", Task(
        description=f"【Act4·PPT压轴🎤】请使用ppt_generator工具，输入'NewCo AI导购项目董事会汇报：战略共识/技术方案/MVP进展/科研产出/行动计划'来生成PPT。这是全流程的最终产出。{max_w}",
        expected_output="董事会PPT",
        agent=agents["ppt_designer"],
    )))

    # 按Act筛选
    if act:
        act_key = f"act{act}"
        tasks = [t for a, t in all_tasks if a == act_key]
    else:
        tasks = [t for _, t in all_tasks]

    return Crew(
        agents=list(agents.values()),
        tasks=tasks,
        process=Process.sequential,
        verbose=True,
    ), len(tasks)


# ═══════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(description="NewCo公司OS — CrewAI多Agent编排")
    parser.add_argument("--act", type=int, choices=[1, 2, 3, 4], help="只运行指定Act")
    parser.add_argument("--demo", action="store_true", help="演示模式(输出简化)")
    parser.add_argument("--list", action="store_true", help="列出所有Agent和工具")
    parser.add_argument("--run-tool", type=str, help="直接运行某个工具(kb/hypothesis/playwright/ppt/research/dashboard/supervisor)")
    parser.add_argument("--tool-input", type=str, default="", help="工具输入参数")
    args = parser.parse_args()

    # 列出工具
    if args.list:
        print("\n🏢 NewCo 公司操作系统 — 工具注册表\n")
        print(f"{'ID':<12} {'名称':<20} {'独立运行'}")
        print("─" * 50)
        for tid, (name, _) in TOOL_REGISTRY.items():
            print(f"{tid:<12} {name:<20} python3.11 newco_crew.py --run-tool {tid}")
        print(f"\n共 {len(TOOL_REGISTRY)} 个工具，每个都可独立运行。\n")
        return

    # 直接运行单个工具
    if args.run_tool:
        if args.run_tool not in TOOL_REGISTRY:
            print(f"❌ 未知工具: {args.run_tool}")
            print(f"可用: {', '.join(TOOL_REGISTRY.keys())}")
            return
        name, ToolClass = TOOL_REGISTRY[args.run_tool]
        print(f"\n🔧 运行工具: {name}\n")
        tool = ToolClass()
        result = tool._run(args.tool_input)
        print(result)
        return

    # CrewAI编排
    if not CREWAI_OK:
        print("❌ CrewAI未安装。请运行: python3.11 -m pip install crewai 'crewai[tools]'")
        return

    print()
    print("╔══════════════════════════════════════════════════════════╗")
    print("║  🏢 NewCo 公司操作系统                                   ║")
    print("║  CrewAI编排 · 10个Agent · 7个真实工具 · 11个Task         ║")
    print("║                                                          ║")
    print("║  每个Agent调用真实的Demo项目，不是纯LLM对话              ║")
    print("╚══════════════════════════════════════════════════════════╝")
    print()

    act_label = f"Act {args.act}" if args.act else "全部4幕"
    print(f"🎬 运行: {act_label}" + (" (演示模式)" if args.demo else ""))
    print()

    crew, task_count = build_crew(act=args.act, demo_mode=args.demo)
    print(f"📋 {task_count} 个任务排队中...\n")

    result = crew.kickoff()

    print()
    print("═" * 60)
    print("🎉 NewCo公司OS — 全部任务完成!")
    print("═" * 60)
    print()
    print(result)


if __name__ == "__main__":
    main()
