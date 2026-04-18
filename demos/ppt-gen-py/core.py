"""PPT generator — real python-pptx PPTX bytes with 5 templates and 4 themes.

Templates: cover, agenda, data, split (text + bar chart), summary.
"""
from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt


def _hex(c: str) -> RGBColor:
    c = c.lstrip("#")
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))


THEMES: dict[str, dict[str, str]] = {
    "blue-orange": {"bg": "#0b1e4f", "primary": "#f97316", "secondary": "#3b82f6", "text": "#f8fafc", "sub": "#cbd5e1"},
    "red-gold":    {"bg": "#450a0a", "primary": "#fbbf24", "secondary": "#ef4444", "text": "#fef3c7", "sub": "#fde68a"},
    "green-mint":  {"bg": "#022c22", "primary": "#34d399", "secondary": "#10b981", "text": "#f0fdf4", "sub": "#a7f3d0"},
    "purple-pink": {"bg": "#3b0764", "primary": "#f472b6", "secondary": "#a855f7", "text": "#fdf4ff", "sub": "#e9d5ff"},
}

# 16:9 slide size
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _apply_bg(slide, color_hex: str) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = _hex(color_hex)
    # Move to back by re-ordering: we added first, so subsequent shapes stack above.
    # python-pptx doesn't expose z-order directly; adding bg first suffices since
    # this function is always called before content shapes are added.


def _text_box(slide, x, y, w, h, text, *, size=18, color="#ffffff", bold=False, align_center=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    if align_center:
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _hex(color)
    return tb


def _render_cover(slide, theme, fields):
    _apply_bg(slide, theme["bg"])
    _text_box(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.4),
              fields.get("title", "Title"), size=48, color=theme["text"], bold=True)
    _text_box(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.8),
              fields.get("subtitle", ""), size=22, color=theme["sub"])
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.9), Inches(1.0), Inches(0.08))
    bar.line.fill.background(); bar.fill.solid(); bar.fill.fore_color.rgb = _hex(theme["primary"])
    _text_box(slide, Inches(0.8), Inches(6.7), Inches(6), Inches(0.5),
              fields.get("company", ""), size=12, color=theme["sub"])
    _text_box(slide, Inches(7.3), Inches(6.7), Inches(5), Inches(0.5),
              fields.get("date", ""), size=12, color=theme["sub"])


def _render_agenda(slide, theme, fields):
    _apply_bg(slide, theme["bg"])
    _text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(1),
              fields.get("title", "Agenda"), size=32, color=theme["text"], bold=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(0.8), Inches(0.06))
    bar.line.fill.background(); bar.fill.solid(); bar.fill.fore_color.rgb = _hex(theme["primary"])
    items = fields.get("items") or (fields.get("content", "") or "").split("\n")
    items = [i for i in items if str(i).strip()]
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11), Inches(5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"{i+1:02d}   {it}"
        run.font.size = Pt(20)
        run.font.color.rgb = _hex(theme["text"])
        p.space_after = Pt(12)


def _render_data(slide, theme, fields):
    _apply_bg(slide, theme["bg"])
    _text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
              fields.get("title", ""), size=20, color=theme["sub"])
    _text_box(slide, Inches(0.8), Inches(1.8), Inches(6), Inches(3.5),
              str(fields.get("big_number", fields.get("number", "78%"))),
              size=120, color=theme["primary"], bold=True)
    _text_box(slide, Inches(7), Inches(2.5), Inches(5.8), Inches(3),
              fields.get("description", ""), size=20, color=theme["text"])
    _text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
              f"Source: {fields.get('source', '')}", size=11, color=theme["sub"])


def _render_split(slide, theme, fields):
    _apply_bg(slide, theme["bg"])
    _text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(1),
              fields.get("title", ""), size=28, color=theme["text"], bold=True)
    bullets = fields.get("bullets") or (fields.get("content", "") or "").split("\n")
    bullets = [b for b in bullets if str(b).strip()][:6]
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.8), Inches(5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = f"▸  {b}"
        r.font.size = Pt(16); r.font.color.rgb = _hex(theme["text"])
        p.space_after = Pt(8)
    # Bar chart on right
    data = fields.get("chart") or {"categories": ["Q1", "Q2", "Q3", "Q4", "Q5"], "values": [65, 48, 82, 34, 71]}
    chart_data = CategoryChartData()
    chart_data.categories = data["categories"]
    chart_data.add_series(fields.get("series_name", "Value"), data["values"])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7), Inches(2.0), Inches(5.8), Inches(4.8), chart_data
    ).chart
    chart.has_title = False
    chart.has_legend = False
    # Color first series
    ser = chart.series[0]
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = _hex(theme["primary"])


def _render_summary(slide, theme, fields):
    _apply_bg(slide, theme["bg"])
    _text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(1),
              fields.get("title", "Summary"), size=28, color=theme["text"], bold=True)
    items = fields.get("columns") or (fields.get("content", "") or "").split("\n")
    items = [i for i in items if str(i).strip()][:3]
    while len(items) < 3:
        items.append("Key takeaway")
    col_w = Inches(3.8); gap = Inches(0.25); x0 = Inches(0.8); y = Inches(2.0); h = Inches(3.8)
    for i, txt in enumerate(items):
        x = Emu(x0 + (col_w + gap) * i)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, h)
        box.line.color.rgb = _hex(theme["primary"])
        box.fill.solid(); box.fill.fore_color.rgb = _hex("#1e293b")
        _text_box(slide, x + Inches(0.3), y + Inches(0.3), col_w - Inches(0.6), Inches(0.8),
                  f"{i+1:02d}", size=32, color=theme["primary"], bold=True)
        _text_box(slide, x + Inches(0.3), y + Inches(1.3), col_w - Inches(0.6), Inches(2.2),
                  txt, size=14, color=theme["text"])
    _text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
              f"{fields.get('company', '')} · {fields.get('date', '')}", size=11, color=theme["sub"])


_RENDERERS = {
    "cover": _render_cover,
    "agenda": _render_agenda,
    "data": _render_data,
    "split": _render_split,
    "summary": _render_summary,
}


def generate_pptx(
    title: str,
    subtitle: str,
    slides: list[dict[str, Any]],
    theme: str = "blue-orange",
) -> bytes:
    """Build a PPTX with a cover auto-prepended if not already in slides."""
    t = THEMES.get(theme, THEMES["blue-orange"])
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # Auto-cover if missing
    if not slides or slides[0].get("template") != "cover":
        _RENDERERS["cover"](prs.slides.add_slide(blank), t,
                            {"title": title, "subtitle": subtitle})
    for s in slides:
        tmpl = s.get("template", "agenda")
        fields = dict(s.get("fields", {}))
        fields.setdefault("title", title)
        renderer = _RENDERERS.get(tmpl, _render_agenda)
        renderer(prs.slides.add_slide(blank), t, fields)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
